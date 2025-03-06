import os
import logging
import hashlib
from dotenv import load_dotenv
from pinecone import Pinecone
import fitz  # PyMuPDF for PDFs
import docx
import pandas as pd
from pptx import Presentation
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_pinecone import Pinecone as PineconeVectorStore
from langchain.text_splitter import RecursiveCharacterTextSplitter
import numpy as np  # Import NumPy to handle NaN values

# ✅ Load environment variables
load_dotenv()

# ✅ Initialize Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ✅ Initialize Pinecone
pinecone_api_key = os.getenv("PINECONE_API_KEY")
index_name = os.getenv("PINECONE_INDEX_NAME")

pc = Pinecone(api_key=pinecone_api_key)

# ✅ Check if Pinecone index exists, if not, create it
if index_name not in [idx["name"] for idx in pc.list_indexes()]:
    logger.info(f"🛠 Creating Pinecone index: {index_name}")
    pc.create_index(
        name=index_name,
        dimension=3072,  # Matches OpenAI embeddings
        metric="dotproduct"
    )

# ✅ Connect to Pinecone Index
index = pc.Index(index_name)

# ✅ Initialize OpenAI Embeddings
embedding_model = OpenAIEmbeddings(model=os.getenv("OPENAI_EMBEDDING_MODEL"))

# ✅ Set up Pinecone VectorStore with namespace support
vector_store = PineconeVectorStore(index=index, embedding=embedding_model, namespace="default")

logger.info("✅ Pinecone & OpenAI Embeddings Initialized Successfully!")

def extract_text(file_path: str) -> str:
    """Extracts text from PDFs, DOCX, PPTX, and XLSX files with better handling."""
    try:
        if file_path.endswith(".pdf"):
            text = ""
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text("text") + "\n"
            return text.strip() if text.strip() else None

        elif file_path.endswith(".docx"):
            doc = docx.Document(file_path)
            text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
            return text if text.strip() else None

        elif file_path.endswith(".xlsx"):
            try:
                df_dict = pd.read_excel(file_path, sheet_name=None, dtype=str, engine="openpyxl")
                text = ""
                for sheet_name, df in df_dict.items():
                    if df.empty:
                        logger.warning(f"⚠️ Empty sheet detected: {sheet_name}")
                        continue  # Skip empty sheets
                    df.dropna(axis=1, how="all", inplace=True)  # Drop empty columns
                    df.dropna(how="all", inplace=True)  # Drop empty rows
                    # ✅ Replace NaN values with empty strings
                    df.fillna("", inplace=True)  # Replace NaN with empty string
                    text += f"\n\nSheet Name: {sheet_name}\n"
                    text += df.to_string(index=False, header=True)
                return text if text.strip() else None
            except Exception as e:
                logger.error(f"❌ Error reading Excel file: {file_path}, Error: {e}")
                return None

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            text = "\n".join(
                [shape.text.strip() for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            )
            return text if text.strip() else None

        else:
            raise ValueError(f"❌ Unsupported file type: {file_path}")

    except Exception as e:
        logger.error(f"❌ Error extracting text from {file_path}: {e}")
        return None

def document_exists_in_pinecone(filename: str) -> bool:
    """Checks if embeddings for a file exist in Pinecone by fetching metadata."""
    try:
        response = index.fetch(ids=[filename])  # Fetch based on filename as the ID
        if response and response.get("vectors"):
            logger.info(f"✅ File '{filename}' already exists in Pinecone. Skipping processing.")
            return True  
        logger.info(f"📄 File '{filename}' not found in Pinecone. Needs processing.")
        return False  
    except Exception as e:
        logger.error(f"❌ Error checking Pinecone for {filename}: {e}")
        return False  

def generate_chunk_id(filename: str, chunk_text: str) -> str:
    """Generates a unique hash ID for each chunk based on filename and chunk text."""
    return hashlib.sha256(f"{filename}-{chunk_text}".encode()).hexdigest()

def clean_context(text: str) -> str:
    """Cleans extracted text by removing 'NaN' and redundant spaces."""
    if not text:
        return ""
    cleaned_text = text.replace("NaN", "").replace("nan", "").strip()
    cleaned_text = " ".join(cleaned_text.split())  # Reduce multiple spaces
    return cleaned_text if cleaned_text else "No relevant information available."

def process_and_store_embeddings(file_path: str, filename: str) -> int:
    """Extracts text, chunks it efficiently, and stores embeddings in Pinecone."""
    
    # ✅ Check if file exists before processing
    if document_exists_in_pinecone(filename):
        logger.info(f"⏭️ Skipping {filename}, embeddings already exist in Pinecone.")
        return 0  

    text = extract_text(file_path)

    if not text:
        logger.warning(f"⚠️ Skipping {filename} due to empty content")
        return 0  

    # ✅ Clean text before processing
    text = clean_context(text)

    # ✅ Log extracted text preview
    logger.info(f"📑 Extracted text preview (first 500 chars): {text[:500]}")

    # ✅ Improved Chunking: Smaller Chunks with Overlap for Better Retrieval
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=512,  
        chunk_overlap=128,  
        separators=["\n\n", "\n", " ", ""]
    )

    chunks = text_splitter.split_text(text)

    if not chunks:
        logger.warning(f"⚠️ No valid text chunks extracted from {filename}")
        return 0  

    # ✅ Generate Unique IDs and Metadata
    chunk_ids = [generate_chunk_id(filename, chunk) for chunk in chunks]
    metadata = [{"source": filename, "chunk_id": chunk_id, "chunk_number": i} for i, chunk_id in enumerate(chunk_ids)]

    # ✅ Log metadata before storing
    logger.info(f"📥 Storing {len(chunks)} chunks in Pinecone for {filename} with metadata: {metadata[:2]}...")  # Log first 2 entries

    # ✅ Store embeddings in Pinecone with metadata
    vector_store.add_texts(chunks, metadatas=metadata, ids=chunk_ids)

    logger.info(f"✅ Successfully stored {len(chunks)} chunks for {filename} in Pinecone")

    return len(chunks)

def get_vector_store():
    """Returns the initialized Pinecone vector store for retrieval."""
    return vector_store
